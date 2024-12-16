from flask import Flask, request, jsonify, send_file, render_template, redirect, url_for
import openai
import os
from PyPDF2 import PdfReader
from docx import Document
from dotenv import load_dotenv
import pandas as pd
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
from datetime import datetime, timedelta
import time
import random
from flask import send_from_directory 

# Carregar variáveis de ambiente do arquivo .env
load_dotenv()

# Configurar a chave da API OpenAI
openai.api_key = os.getenv("OPENAI_API_KEY")

app = Flask(__name__)

# Aprendizado prévio incorporado
LEARNING_DATA_PATHS = [
    "data/Proposta-Comercial.pdf",
    "data/CLIENTE_Escopo_acordado.pdf",
    "data/CLIENTE_Kick-off.pdf"
]

def load_learning_data():
    learning_data = []
    for file_path in LEARNING_DATA_PATHS:
        # Leia o texto do PDF
        pdf_reader = PdfReader(file_path)
        text = "".join([page.extract_text() for page in pdf_reader.pages])
        learning_data.append(text)
    return learning_data

@app.route('/')
def upload_page():
    return render_template('upload.html')

@app.route('/generate-scope', methods=['POST'])
def generate_scope():
    # Verifique se o arquivo foi enviado
    if 'file' not in request.files:
        return jsonify({"error": "Nenhum arquivo enviado"}), 400

    # Carregar o arquivo PDF enviado pelo usuário
    file = request.files['file']
    pdf_reader = PdfReader(file)
    current_text = "".join([page.extract_text() for page in pdf_reader.pages])

    # Capturar os campos do formulário
    client_name = request.form.get('client_name')
    project_type = request.form.get('project_type')
    # Debug para verificar os valores capturados
    print(f"Cliente: {client_name}")
    print(f"Tipo de Projeto: {project_type}")

    # Carregar dados de aprendizado
    learning_examples = load_learning_data()

    # Construir o prompt com exemplos e a proposta atual
    prompt = "Você é um especialista em criar escopos acordados. Aqui estão exemplos prévios:\n\n"
    for example in learning_examples:
        prompt += f"Exemplo de Escopo Acordado:\n{example}\n\n"
    prompt += f"Proposta Atual:\n{current_text}\n\nCrie o escopo baseado nisso."

    # Enviar o prompt para a OpenAI
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "Você é um especialista em criar escopos acordados."},
            {"role": "user", "content": prompt}
        ],
        max_tokens=2000,
        temperature=0.7
    )

    # Obter o texto gerado
    generated_text = response["choices"][0]["message"]["content"]


    # Criar o documento Word
    document = Document()
    document.add_heading('Escopo Acordado', level=1)
    document.add_paragraph(generated_text)
    word_path = "data/escopo_acordado.docx"
    document.save(word_path)

    # Gerar cronograma dinâmico
    entregas = []
    tarefas = []

    # Identificar entregas no escopo
    for line in generated_text.splitlines():
        if line.strip().startswith("•") or line.strip().startswith("-"):
            entregas.append(line.strip().lstrip("•- ").strip())

    # Criar tarefas para cada entrega
    for entrega in entregas:
        tarefas.append({
            "entrega": entrega,
            "tarefas": [f"Planejar {entrega}", f"Executar {entrega}", f"Finalizar {entrega}"]
        })

    # Gerar datas
    data_atual = datetime.now()
    gantt_data = []
    for tarefa in tarefas:
        for i, subtarefa in enumerate(tarefa["tarefas"]):
            inicio = data_atual + timedelta(days=i * 7)  # Cada subtarefa começa 7 dias após a anterior
            fim = inicio + timedelta(days=6)  # Cada subtarefa dura 7 dias
            gantt_data.append({
                "Tarefa": subtarefa,
                "Início": inicio.strftime("%Y-%m-%d"),
                "Fim": fim.strftime("%Y-%m-%d")
            })

    df = pd.DataFrame(gantt_data)
    excel_path = "data/cronograma.xlsx"
    df.to_excel(excel_path, index=False)

    # Criar a apresentação em PowerPoint
    ppt = Presentation("data\modelo_kickoff.pptx")
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "<TITULO_ESC>" in shape.text:
                    shape.text = "Resumo do Escopo"
                if "<RESUMO_ESC>" in shape.text:
                    shape.text = generated_text[:500]  # Resumo do escopo
                if "<CRONOGRAMA>" in shape.text:
                    shape.text = "Cronograma será anexado na apresentação."

    ppt_path = "data/apresentacao_kickoff.pptx"
    ppt.save(ppt_path)

    return render_template(
        'result.html',
        word_path=word_path,
        excel_path=excel_path,
        ppt_path=ppt_path,
        generated_text=generated_text,  # Passar o texto gerado
        client_name="Cliente X",
        project_type="Tipo de Projeto"
    )

@app.route('/data/<path:filename>')
def serve_file(filename):
    directory = os.path.join(os.getcwd(), "data")  # Caminho absoluto para a pasta "data"
    if not os.path.exists(os.path.join(directory, filename)):
        return f"Erro: O arquivo {filename} não foi encontrado no diretório {directory}.", 404
    return send_from_directory(directory, filename)


@app.route('/init-immersion', methods=['POST'])
def init_immersion():
    client_name = request.form.get('client_name')
    project_type = request.form.get('project_type')

    print(f"Cliente: {client_name}")
    print(f"Tipo de Projeto: {project_type}")

    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36",
        "Accept-Language": "en-US,en;q=0.9",
    }
    
    news_links = []
    glassdoor_data = []
    error_messages = []  # Lista para mensagens de erro

    try:
        # Raspagem CNN Brasil
        cnn_url = f"https://www.cnnbrasil.com.br/?s=%22{client_name}%22&orderby=date&order=desc"
        time.sleep(random.uniform(2, 5))
        cnn_response = requests.get(cnn_url, headers=headers, timeout=10)
        cnn_soup = BeautifulSoup(cnn_response.text, 'html.parser')
        for item in cnn_soup.find_all("h3", class_="news-item-header__title"):
            news_links.append({"title": item.get_text(), "url": cnn_url})
    except Exception as e:
        error_messages.append("Não foi possível obter notícias da CNN Brasil.")

    try:
        # Raspagem Glassdoor
        glassdoor_url = f"https://www.glassdoor.com.br/Vis%C3%A3o-geral/Trabalhar-na-{client_name.replace(' ', '-')}-EI_IE9323739.13,27.htm"
        time.sleep(random.uniform(2, 5))
        glassdoor_response = requests.get(glassdoor_url, headers=headers, timeout=10)
        glassdoor_soup = BeautifulSoup(glassdoor_response.text, 'html.parser')

        for script_tag in glassdoor_soup.find_all('script', type="application/ld+json"):
            try:
                data = json.loads(script_tag.string)
                if isinstance(data, dict) and "@type" in data and data["@type"] == "Question":
                    glassdoor_data.append({
                        "question": data["name"],
                        "answer": data["acceptedAnswer"]["text"]
                    })
            except (json.JSONDecodeError, KeyError):
                continue
    except Exception as e:
        error_messages.append("Não foi possível obter dados do Glassdoor.")

    # Processar análise com ChatGPT
    glassdoor_text = "\n".join(
        [f"Pergunta: {item['question']}\nResposta: {item['answer']}" for item in glassdoor_data]
    )
    prompt = (
        f"Você é um especialista em análise de marca empregadora e cultura organizacional. "
        f"Baseado nas informações coletadas no Glassdoor sobre a empresa {client_name}, elabore um resumo "
        f"detalhado, destacando os pontos fortes, fracos e qualquer percepção geral que possa ser extraída dos dados.\n\n"
        f"Dados do Glassdoor:\n{glassdoor_text}"
    )

    analysis = "Análise não disponível."
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "Você é um analista de marca empregadora."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=2000,
            temperature=0.7
        )
        analysis = response["choices"][0]["message"]["content"]
    except Exception as e:
        error_messages.append("Não foi possível processar a análise da OpenAI.")

    # Renderizar a página com as informações ou mensagens de erro
    return render_template(
        'immersion.html',
        client_name=client_name,
        project_type=project_type,
        news_links=news_links,
        glassdoor_data=glassdoor_data,
        analysis=analysis,
        error_messages=error_messages
    )



if __name__ == "__main__":
    app.run(debug=True)